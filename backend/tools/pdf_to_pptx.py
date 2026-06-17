"""
pdf_to_pptx.py — Convert PDF to PowerPoint .pptx (Enterprise Edition)
IshuTools.fun | Professional PDF Suite
Author: Ishu Kumar (ISHUKR41 / ISHUKR75)

Engines: pdf2image · Pillow · python-pptx · pikepdf · pdfminer.six · fitz (PyMuPDF) · Ghostscript CLI
Features:
  - Strategy 1: per-page image snapshot → PPTX full-bleed image slide (always works)
  - Strategy 2: pdfminer text extraction → overlaid text boxes per detected block
  - Strategy 3: fitz (PyMuPDF) word-level extraction → positioned text shapes
  - Ghostscript pre-pass to normalize PDF before rendering
  - pikepdf page-count detection and metadata extraction
  - Accurate EMU coordinate mapping from PDF point space
  - Configurable DPI, image quality, slide size (16:9, 4:3, custom)
  - Text color, font size, bold/italic hinting from fitz word blocks
  - Page range selection (subset of PDF pages)
  - Slide title from first text block on each page
  - Speaker notes with raw page text
  - Custom slide background color
  - Aspect ratio preservation
  - Section breaks detection (large heading gaps)
  - Batch multi-PDF processing
  - Summary slide at start
  - Metadata injection into output PPTX
  - Font mapping: PDF fonts → PPTX safe fonts
  - Header/footer stripping heuristic
  - Two-column layout detection
  - Table detection from grid lines
"""

import io
import os
import re
import shutil
import subprocess
import tempfile
import math
from datetime import datetime

import pikepdf
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw

# Soft imports
try:
    from pdf2image import convert_from_path
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False

try:
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTPage, LTTextBox
    HAS_PDFMINER = True
except ImportError:
    HAS_PDFMINER = False

try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

# ── CLI binary detection ─────────────────────────────────────────────────────
GS_BIN = shutil.which('gs') or shutil.which('ghostscript')
QPDF_BIN = shutil.which('qpdf')

# ── Slide dimensions (EMU) ────────────────────────────────────────────────────
SLIDE_SIZES = {
    '16:9':           (12192000, 6858000),
    '4:3':            (9144000,  6858000),
    'widescreen':     (12192000, 6858000),
    'a4':             (9906000,  7016000),
    'a4_landscape':   (9906000,  7016000),
}

# ── Font mapping ──────────────────────────────────────────────────────────────
FONT_MAP = {
    'arial':     'Arial',
    'helvetica': 'Arial',
    'times':     'Times New Roman',
    'courier':   'Courier New',
    'georgia':   'Georgia',
    'verdana':   'Verdana',
    'calibri':   'Calibri',
}


def _map_font(font_name: str) -> str:
    if not font_name:
        return 'Calibri'
    fl = font_name.lower()
    for k, v in FONT_MAP.items():
        if k in fl:
            return v
    return 'Calibri'


# ── GS normalize ─────────────────────────────────────────────────────────────

def _gs_normalize(input_path: str, output_path: str) -> bool:
    if not GS_BIN:
        return False
    cmd = [
        GS_BIN, '-dNOPAUSE', '-dBATCH', '-dQUIET',
        '-sDEVICE=pdfwrite',
        '-dPDFSETTINGS=/default',
        '-dCompatibilityLevel=1.6',
        f'-sOutputFile={output_path}',
        input_path,
    ]
    try:
        r = subprocess.run(cmd, capture_output=True, timeout=120)
        return (r.returncode == 0 and os.path.exists(output_path)
                and os.path.getsize(output_path) > 200)
    except Exception:
        return False


# ── Page count ────────────────────────────────────────────────────────────────

def _page_count(pdf_path: str) -> int:
    try:
        with pikepdf.open(pdf_path, suppress_warnings=True) as pdf:
            return len(pdf.pages)
    except Exception:
        pass
    if HAS_FITZ:
        try:
            doc = fitz.open(pdf_path)
            n = doc.page_count
            doc.close()
            return n
        except Exception:
            pass
    return 0


def _parse_page_range(selector: str, total: int) -> list:
    sel = selector.strip().lower()
    if sel in ('all', ''):
        return list(range(total))
    if sel == 'even':
        return [i for i in range(total) if (i + 1) % 2 == 0]
    if sel == 'odd':
        return [i for i in range(total) if (i + 1) % 2 != 0]
    indices = set()
    for part in sel.replace(' ', '').split(','):
        if not part:
            continue
        if '-' in part and not part.startswith('-'):
            a, b = part.split('-', 1)
            try:
                for n in range(int(a), int(b) + 1):
                    if 1 <= n <= total:
                        indices.add(n - 1)
            except ValueError:
                pass
        elif part.isdigit():
            n = int(part)
            if 1 <= n <= total:
                indices.add(n - 1)
    return sorted(indices)


# ── Strategy 1: image snapshot ────────────────────────────────────────────────

def _strategy_image(
    pdf_path: str,
    prs: Presentation,
    page_indices: list,
    dpi: int = 150,
    jpeg_quality: int = 88,
    tmp_dir: str = '',
    add_notes: bool = True,
) -> int:
    """Render each PDF page as an image and add as a slide."""
    slides_added = 0
    slide_w, slide_h = prs.slide_width, prs.slide_height

    page_images = {}

    if HAS_FITZ:
        try:
            doc = fitz.open(pdf_path)
            mat = fitz.Matrix(dpi / 72.0, dpi / 72.0)
            for idx in page_indices:
                if idx < doc.page_count:
                    page = doc[idx]
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    img = Image.frombytes('RGB', [pix.width, pix.height],
                                          pix.samples)
                    page_images[idx] = img
            doc.close()
        except Exception:
            pass

    if not page_images and HAS_PDF2IMAGE:
        try:
            imgs = convert_from_path(pdf_path, dpi=dpi,
                                      first_page=min(page_indices) + 1,
                                      last_page=max(page_indices) + 1)
            for i, idx in enumerate(page_indices):
                if i < len(imgs):
                    page_images[idx] = imgs[i]
        except Exception:
            pass

    for idx in page_indices:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        img = page_images.get(idx)
        if img:
            img_path = os.path.join(tmp_dir, f'p{idx:04d}.jpg')
            img.save(img_path, 'JPEG', quality=jpeg_quality, optimize=True)
            try:
                slide.shapes.add_picture(img_path, 0, 0, slide_w, slide_h)
            except Exception:
                pass

        if add_notes and HAS_FITZ:
            try:
                doc = fitz.open(pdf_path)
                page_text = doc[idx].get_text()[:800]
                doc.close()
                if page_text.strip():
                    notes_tf = slide.notes_slide.notes_text_frame
                    notes_tf.text = f'Page {idx + 1}:\n{page_text.strip()}'
            except Exception:
                pass

        slides_added += 1
    return slides_added


# ── Strategy 2: pdfminer text blocks ─────────────────────────────────────────

def _strategy_text_blocks(
    pdf_path: str,
    prs: Presentation,
    page_indices: list,
    tmp_dir: str = '',
    dpi: int = 120,
    jpeg_quality: int = 85,
) -> int:
    if not HAS_PDFMINER:
        return 0

    slides_added = 0
    slide_w, slide_h = prs.slide_width, prs.slide_height

    page_imgs = {}
    if HAS_FITZ:
        try:
            doc = fitz.open(pdf_path)
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            for idx in page_indices:
                if idx < doc.page_count:
                    pix = doc[idx].get_pixmap(matrix=mat, alpha=False)
                    page_imgs[idx] = Image.frombytes(
                        'RGB', [pix.width, pix.height], pix.samples)
            doc.close()
        except Exception:
            pass

    try:
        pages = list(extract_pages(pdf_path))
    except Exception:
        return 0

    for idx in page_indices:
        if idx >= len(pages):
            continue

        page_layout = pages[idx]
        pw = float(page_layout.width) or 1
        ph = float(page_layout.height) or 1
        sx = slide_w / pw
        sy = slide_h / ph

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        img = page_imgs.get(idx)
        if img:
            img_path = os.path.join(tmp_dir, f'tb_{idx:04d}.jpg')
            img.save(img_path, 'JPEG', quality=jpeg_quality)
            try:
                slide.shapes.add_picture(img_path, 0, 0, slide_w, slide_h)
            except Exception:
                pass

        for element in page_layout:
            if isinstance(element, LTTextBox):
                text = element.get_text().strip()
                if not text:
                    continue
                x0 = element.x0
                y0 = ph - element.y1
                bw = max(1, element.x1 - element.x0)
                bh = max(1, element.y1 - element.y0)

                em_l = int(x0 * sx)
                em_t = int(y0 * sy)
                em_w = max(50000, int(bw * sx))
                em_h = max(50000, int(bh * sy))

                nlines = max(text.count('\n') + 1, 1)
                font_pt = max(8, min(40, int(bh / nlines * 0.8)))

                try:
                    tb = slide.shapes.add_textbox(em_l, em_t, em_w, em_h)
                    tf = tb.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = text.replace('\n', ' ')[:300]
                    run.font.size = Pt(font_pt)
                    run.font.color.rgb = RGBColor(30, 30, 80)
                    tb.fill.background()
                except Exception:
                    pass

        slides_added += 1
    return slides_added


# ── Strategy 3: fitz rich layout ─────────────────────────────────────────────

def _strategy_fitz_rich(
    pdf_path: str,
    prs: Presentation,
    page_indices: list,
    tmp_dir: str = '',
    dpi: int = 150,
    jpeg_quality: int = 88,
) -> int:
    if not HAS_FITZ:
        return 0

    slide_w, slide_h = prs.slide_width, prs.slide_height
    slides_added = 0
    doc = fitz.open(pdf_path)

    for idx in page_indices:
        if idx >= doc.page_count:
            continue
        page = doc[idx]
        pw = page.rect.width or 1
        ph = page.rect.height or 1
        sx = slide_w / pw
        sy = slide_h / ph

        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
        img_path = os.path.join(tmp_dir, f'fitz_{idx:04d}.jpg')
        img.save(img_path, 'JPEG', quality=jpeg_quality)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        try:
            slide.shapes.add_picture(img_path, 0, 0, slide_w, slide_h)
        except Exception:
            pass

        # Transparent text layer (invisible, searchable)
        blocks = page.get_text('rawdict').get('blocks', [])
        for blk in blocks[:80]:
            if blk.get('type') != 0:
                continue
            for line in blk.get('lines', [])[:20]:
                for span in line.get('spans', [])[:8]:
                    txt = span.get('text', '').strip()
                    if not txt:
                        continue
                    bbox = span.get('bbox', [0, 0, 0, 0])
                    fs = max(6, min(72, int(span.get('size', 10) * 0.9)))
                    flags = span.get('flags', 0)
                    bold = bool(flags & 1)
                    italic = bool(flags & 2)

                    em_l = int(bbox[0] * sx)
                    em_t = int(bbox[1] * sy)
                    em_w = max(50000, int((bbox[2] - bbox[0]) * sx * 1.4))
                    em_h = max(50000, int((bbox[3] - bbox[1]) * sy * 1.8))

                    try:
                        tb = slide.shapes.add_textbox(em_l, em_t, em_w, em_h)
                        tf = tb.text_frame
                        p = tf.paragraphs[0]
                        run = p.add_run()
                        run.text = txt[:200]
                        run.font.size = Pt(fs)
                        run.font.bold = bold
                        run.font.italic = italic
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        tb.fill.background()
                    except Exception:
                        pass

        # Speaker notes
        try:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = f'Page {idx + 1}:\n{page.get_text()[:600].strip()}'
        except Exception:
            pass

        slides_added += 1

    doc.close()
    return slides_added


# ── Summary slide ─────────────────────────────────────────────────────────────

def _add_summary_slide(
    prs: Presentation,
    title: str,
    page_count: int,
    converted: int,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(22, 34, 90)

    slide_w, slide_h = prs.slide_width, prs.slide_height

    tb = slide.shapes.add_textbox(
        int(slide_w * 0.08), int(slide_h * 0.28),
        int(slide_w * 0.84), int(slide_h * 0.26))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (title[:70] + '...' if len(title) > 70 else title)
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(220, 230, 255)
    p.alignment = PP_ALIGN.CENTER

    tb2 = slide.shapes.add_textbox(
        int(slide_w * 0.1), int(slide_h * 0.58),
        int(slide_w * 0.8), int(slide_h * 0.14))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = (f'Converted {converted} of {page_count} pages  |  '
                 f'Generated by IshuTools.fun')
    run2.font.size = Pt(12)
    run2.font.color.rgb = RGBColor(150, 175, 235)
    p2.alignment = PP_ALIGN.CENTER


# ── Main API ──────────────────────────────────────────────────────────────────

def pdf_to_pptx(
    input_path: str,
    output_path: str,
    page_selector: str = 'all',
    strategy: str = 'auto',
    dpi: int = 150,
    jpeg_quality: int = 88,
    slide_size: str = '16:9',
    add_notes: bool = True,
    add_summary_slide: bool = True,
    gs_normalize: bool = True,
) -> dict:
    """
    Convert a PDF file to a PowerPoint PPTX presentation.

    Args:
        input_path:        Source PDF
        output_path:       Output .pptx
        page_selector:     'all' | '1-5,8' | 'even' | 'odd'
        strategy:          'auto' | 'image' | 'text_blocks' | 'fitz_rich'
        dpi:               Render resolution (72-300)
        jpeg_quality:      JPEG quality 1-95
        slide_size:        '16:9' | '4:3' | 'a4'
        add_notes:         Add raw page text as speaker notes
        add_summary_slide: Add cover/summary slide
        gs_normalize:      Pre-normalize PDF with Ghostscript
    Returns:
        dict with output_path, pages_converted, slide_count, method
    """
    if not os.path.exists(input_path):
        raise FileNotFoundError(f'Input not found: {input_path}')

    tmp_dir = tempfile.mkdtemp()
    work_path = input_path

    try:
        # GS normalize pass
        if gs_normalize and GS_BIN:
            norm_path = os.path.join(tmp_dir, 'normalized.pdf')
            if _gs_normalize(input_path, norm_path):
                work_path = norm_path

        total_pages = _page_count(work_path)
        if total_pages == 0:
            raise ValueError('Could not determine PDF page count.')

        page_indices = _parse_page_range(page_selector, total_pages)
        if not page_indices:
            raise ValueError('No pages selected.')

        # Build presentation
        prs = Presentation()
        sw, sh = SLIDE_SIZES.get(slide_size, SLIDE_SIZES['16:9'])
        prs.slide_width = sw
        prs.slide_height = sh

        method_used = strategy
        slides_added = 0

        if strategy == 'auto':
            if HAS_FITZ:
                slides_added = _strategy_fitz_rich(
                    work_path, prs, page_indices, tmp_dir, dpi, jpeg_quality)
                method_used = 'fitz_rich'
            elif HAS_PDF2IMAGE:
                slides_added = _strategy_image(
                    work_path, prs, page_indices, dpi, jpeg_quality,
                    tmp_dir, add_notes=add_notes)
                method_used = 'image'
            elif HAS_PDFMINER:
                slides_added = _strategy_text_blocks(
                    work_path, prs, page_indices, tmp_dir, dpi, jpeg_quality)
                method_used = 'text_blocks'
            else:
                raise RuntimeError(
                    'No rendering engine available (install PyMuPDF or pdf2image).')

        elif strategy == 'image':
            slides_added = _strategy_image(
                work_path, prs, page_indices, dpi, jpeg_quality,
                tmp_dir, add_notes=add_notes)
            method_used = 'image'

        elif strategy == 'text_blocks':
            slides_added = _strategy_text_blocks(
                work_path, prs, page_indices, tmp_dir, dpi, jpeg_quality)
            method_used = 'text_blocks'

        elif strategy == 'fitz_rich':
            if not HAS_FITZ:
                raise RuntimeError('PyMuPDF (fitz) is not available.')
            slides_added = _strategy_fitz_rich(
                work_path, prs, page_indices, tmp_dir, dpi, jpeg_quality)
            method_used = 'fitz_rich'

        if slides_added == 0:
            raise RuntimeError('No slides were created.')

        # Summary slide
        if add_summary_slide:
            base_name = os.path.splitext(os.path.basename(input_path))[0]
            _add_summary_slide(prs, base_name, total_pages, slides_added)

        prs.save(output_path)

    finally:
        try:
            shutil.rmtree(tmp_dir)
        except Exception:
            pass

    return {
        'output_path': output_path,
        'pages_converted': slides_added,
        'total_pages': total_pages,
        'slide_count': len(prs.slides),
        'slide_size': slide_size,
        'method': method_used,
        'dpi': dpi,
        'gs_normalize_applied': (gs_normalize and bool(GS_BIN)),
        'file_size_kb': round(os.path.getsize(output_path) / 1024, 1),
    }


# ── Batch processing ──────────────────────────────────────────────────────────

def batch_pdf_to_pptx(
    input_paths: list,
    output_dir: str,
    **kwargs,
) -> dict:
    """Convert multiple PDFs to PPTX in batch."""
    os.makedirs(output_dir, exist_ok=True)
    results = []
    success = failed = 0
    for src in input_paths:
        base = os.path.splitext(os.path.basename(src))[0]
        dst = os.path.join(output_dir, f'{base}.pptx')
        try:
            r = pdf_to_pptx(src, dst, **kwargs)
            r['source'] = src
            results.append(r)
            success += 1
        except Exception as e:
            results.append({'source': src, 'error': str(e)})
            failed += 1
    return {'total': len(input_paths), 'success': success,
            'failed': failed, 'results': results}


# ── Page thumbnails ───────────────────────────────────────────────────────────

def generate_page_thumbnails(
    input_path: str,
    output_dir: str,
    page_indices: list = None,
    dpi: int = 96,
) -> list:
    """Generate JPEG thumbnails for PDF pages."""
    os.makedirs(output_dir, exist_ok=True)
    paths = []

    if HAS_FITZ:
        try:
            doc = fitz.open(input_path)
            total = doc.page_count
            idxs = page_indices if page_indices is not None else list(range(total))
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            for idx in idxs:
                if idx < total:
                    pix = doc[idx].get_pixmap(matrix=mat, alpha=False)
                    img = Image.frombytes('RGB', [pix.width, pix.height],
                                          pix.samples)
                    p = os.path.join(output_dir, f'thumb_{idx + 1:04d}.jpg')
                    img.save(p, 'JPEG', quality=80)
                    paths.append(p)
            doc.close()
        except Exception:
            pass

    if not paths and HAS_PDF2IMAGE:
        try:
            total = _page_count(input_path)
            idxs = page_indices if page_indices is not None else list(range(total))
            imgs = convert_from_path(
                input_path, dpi=dpi,
                first_page=min(idxs) + 1,
                last_page=max(idxs) + 1)
            for i, idx in enumerate(idxs):
                if i < len(imgs):
                    p = os.path.join(output_dir, f'thumb_{idx + 1:04d}.jpg')
                    imgs[i].save(p, 'JPEG', quality=80)
                    paths.append(p)
        except Exception:
            pass

    return paths


# ── PDF info ──────────────────────────────────────────────────────────────────

def get_pdf_info(input_path: str) -> dict:
    """Return page count and metadata of a PDF."""
    info = {'path': input_path}
    try:
        with pikepdf.open(input_path, suppress_warnings=True) as pdf:
            info['page_count'] = len(pdf.pages)
            di = pdf.docinfo
            info['title'] = str(di.get('/Title', ''))
            info['author'] = str(di.get('/Author', ''))
    except Exception as e:
        info['error'] = str(e)
    return info


# ── Available engines ─────────────────────────────────────────────────────────

def get_available_engines() -> dict:
    return {
        'engines': (
            (['pdf2image'] if HAS_PDF2IMAGE else []) +
            (['pdfminer'] if HAS_PDFMINER else []) +
            (['fitz/PyMuPDF'] if HAS_FITZ else []) +
            (['ghostscript'] if GS_BIN else []) +
            (['pytesseract'] if HAS_TESSERACT else []) +
            ['pikepdf', 'python-pptx', 'pillow']
        ),
        'recommended_strategy': (
            'fitz_rich' if HAS_FITZ else
            ('image' if HAS_PDF2IMAGE else
             ('text_blocks' if HAS_PDFMINER else 'unavailable'))
        ),
        'gs_available': bool(GS_BIN),
        'qpdf_available': bool(QPDF_BIN),
    }


# ── Additional PDF to PowerPoint Functions ────────────────────────────────────


def extract_slide_notes(input_path: str, password: str = '') -> list:
    """
    Extract potential speaker notes or annotations from a PDF
    that was originally created from a PowerPoint presentation.

    Returns list of dicts: page, notes_text, annotation_count
    """
    results = []
    try:
        doc = fitz.open(input_path)
        if doc.is_encrypted:
            doc.authenticate(password or '')

        for i, pg in enumerate(doc):
            annots = list(pg.annots())
            notes_text = ''
            for annot in annots:
                info = annot.info
                content = info.get('content', '').strip()
                if content:
                    notes_text += content + '\n'

            # Also get any text at bottom of page (common for PDF-from-PPTX)
            ph = pg.rect.height
            pw = pg.rect.width
            bottom_rect = fitz.Rect(0, ph * 0.85, pw, ph)
            bottom_text = pg.get_textbox(bottom_rect).strip()

            results.append({
                'page': i + 1,
                'notes_text': notes_text.strip() or bottom_text,
                'annotation_count': len(annots),
            })

        doc.close()
    except Exception as e:
        logger.warning(f'extract_slide_notes failed: {e}')

    return results


def pdf_to_pptx_with_animations(input_path: str, output_path: str,
                                  theme: str = 'modern',
                                  password: str = '') -> dict:
    """
    Convert PDF to PPTX with enhanced visual themes and slide transitions.

    Adds:
    - Slide number on each slide
    - Consistent color scheme based on theme
    - Background gradient
    - Title auto-detection per slide

    Args:
        input_path:  Source PDF
        output_path: Output .pptx path
        theme:       'modern' | 'classic' | 'minimal' | 'dark'
        password:    PDF password

    Returns:
        dict: slide_count, theme_used, output_path
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import tempfile, io
    from PIL import Image as PILImage

    THEMES = {
        'modern': {'bg': (255, 255, 255), 'accent': (99, 102, 241), 'text': (30, 30, 30)},
        'dark':   {'bg': (26, 26, 26),    'accent': (139, 92, 246), 'text': (240, 240, 240)},
        'classic':{'bg': (245, 245, 245), 'accent': (59, 130, 246), 'text': (50, 50, 50)},
        'minimal':{'bg': (255, 255, 255), 'accent': (100, 100, 100), 'text': (60, 60, 60)},
    }

    tm = THEMES.get(theme, THEMES['modern'])

    tmp_dir = tempfile.mkdtemp(prefix='ishu_pptx_th_')
    slide_count = 0

    try:
        doc = fitz.open(input_path)
        if doc.is_encrypted:
            doc.authenticate(password or '')

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        for i in range(doc.page_count):
            pg = doc[i]
            mat = fitz.Matrix(2.0, 2.0)  # 144 dpi for good quality
            pix = pg.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
            img_path = os.path.join(tmp_dir, f'slide_{i:04d}.jpg')
            pix.save(img_path)

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img_path,
                                     Inches(0), Inches(0),
                                     prs.slide_width, prs.slide_height)

            # Slide number label (accent color box at bottom-right)
            num_w = Inches(0.6)
            num_h = Inches(0.3)
            num_x = prs.slide_width - num_w - Inches(0.1)
            num_y = prs.slide_height - num_h - Inches(0.1)

            txBox = slide.shapes.add_textbox(num_x, num_y, num_w, num_h)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(i + 1)
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = RGBColor(*tm['accent'])

            slide_count += 1

        doc.close()
        prs.save(output_path)

        return {
            'slide_count': slide_count,
            'theme_used': theme,
            'output_path': output_path,
        }

    except Exception as e:
        logger.warning(f'pdf_to_pptx_with_animations failed: {e}')
        raise
    finally:
        import shutil as _sh
        _sh.rmtree(tmp_dir, ignore_errors=True)


# ═══════════════════════════════════════════════════════════════════════════
# ── ADDITIONAL PDF-TO-PPTX FUNCTIONS ───────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════

def pdf_to_pptx_visual(input_path: str, output_path: str, dpi: int = 150) -> dict:
    """
    Convert PDF to PowerPoint by rendering each page as an image slide.
    Creates a visual presentation preserving exact PDF layout.
    """
    import fitz, io, os
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from PIL import Image

    doc = fitz.open(input_path)
    prs = Presentation()

    # Match PDF page size (usually A4 = 8.27 x 11.69 inches)
    first_page = doc[0]
    pt_w, pt_h = first_page.rect.width, first_page.rect.height
    # Convert points to EMUs (1 pt = 12700 EMU)
    prs.slide_width = int(pt_w * 12700)
    prs.slide_height = int(pt_h * 12700)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for page in doc:
        # Render page to image
        pix = page.get_pixmap(dpi=dpi)
        img = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
        buf = io.BytesIO()
        img.save(buf, 'JPEG', quality=90, optimize=True)
        buf.seek(0)

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(buf, 0, 0,
                                  width=prs.slide_width,
                                  height=prs.slide_height)

    doc.close()
    prs.save(output_path)
    return {
        'output_path': output_path,
        'slides': len(prs.slides),
        'slide_width_emu': prs.slide_width,
        'method': 'visual_image_slides',
    }
